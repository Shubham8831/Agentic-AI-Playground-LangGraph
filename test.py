"""
Streamlit Advanced PPT Maker (prototype with Groq Cloud)

- Takes a text outline or upload
- Calls Groq Cloud LLM to convert outline -> structured slide JSON
- Renders preview + builds .pptx using python-pptx

Setup:
- pip install streamlit groq python-pptx python-docx
- set GROQ_API_KEY env var

Run:
$ streamlit run streamlit_advanced_ppt_maker.py
"""

import streamlit as st
import os
import json
import tempfile
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt

try:
    from docx import Document
except Exception:
    Document = None

from groq import Groq

GROQ_API_KEY = os.environ.get("GROQ_API_KEY")
client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None

# ------------------------- Helper functions -------------------------

def call_llm_for_slides(outline_text: str, temperature: float = 0.0) -> dict:
    """Call Groq Cloud LLM to convert outline to structured slide JSON."""
    if not client:
        st.error("GROQ_API_KEY not set. Using fallback.")
        return fallback_outline_to_slides(outline_text)

    system_prompt = (
        "You are a helpful assistant that converts an input brief or outline into a slide-by-slide JSON structure. "
        "Return ONLY valid JSON. The JSON must be an object with key 'slides' (array of slide objects). "
        "Each slide object: 'title' (string), 'type' (title, bullet, image, chart, section), "
        "'bullets' (array of strings, optional), 'notes' (string, optional). Keep bullets concise."
    )

    user_prompt = f"Convert the following brief into up to 6 slides. Brief:\n\n{outline_text}\n\nReturn only JSON."

    try:
        resp = client.chat.completions.create(
            model="llama-3.1-70b-versatile",  # or another Groq available model
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=temperature,
            max_tokens=800,
        )
        text = resp.choices[0].message.content.strip()
        json_text = extract_json_block(text)
        return json.loads(json_text)
    except Exception as e:
        st.warning(f"Groq call failed: {e}. Using fallback.")
        return fallback_outline_to_slides(outline_text)


def extract_json_block(text: str) -> str:
    import re
    m = re.search(r"```(?:json)?\s*(\{[\s\S]*\})\s*```", text, flags=re.IGNORECASE)
    if m:
        return m.group(1)
    m2 = re.search(r"(\{[\s\S]*\})", text)
    if m2:
        return m2.group(1)
    return text


def fallback_outline_to_slides(text: str) -> dict:
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    slides = []
    if not lines:
        return {"slides": []}
    slides.append({"title": lines[0], "type": "title", "notes": "Generated title slide"})
    cur = []
    for l in lines[1:]:
        cur.append(l)
        if len(cur) == 3:
            slides.append({"title": cur[0][:40], "type": "bullet", "bullets": cur, "notes": ""})
            cur = []
    if cur:
        slides.append({"title": cur[0][:40], "type": "bullet", "bullets": cur, "notes": ""})
    return {"slides": slides}


def build_pptx_from_json(slide_json: dict) -> str:
    prs = Presentation()
    for s in slide_json.get("slides", []):
        typ = s.get("type", "bullet")
        if typ == "title":
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            if "title" in s:
                slide.shapes.title.text = s.get("title", "")
            if s.get("notes"):
                slide.notes_slide.notes_text_frame.text = s.get("notes")
        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = s.get("title", "")
            body_shape = None
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 1:
                    body_shape = shape
                    break
            if body_shape is None:
                left, top, width, height = Inches(1), Inches(1.8), Inches(8), Inches(4.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
            else:
                tf = body_shape.text_frame
                tf.clear()

            bullets = s.get("bullets", [])
            if bullets:
                for i, b in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = b
                    p.level = 0 if i == 0 else 1
                    p.font.size = Pt(18)
            else:
                p = tf.paragraphs[0]
                p.text = s.get("notes", "")

            if s.get("notes"):
                slide.notes_slide.notes_text_frame.text = s.get("notes")

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name

# ------------------------- Streamlit UI -------------------------

st.set_page_config(page_title="Streamlit PPT Maker", layout="wide")
st.title("Streamlit Advanced PPT Maker — Groq Cloud Prototype")

with st.sidebar:
    st.header("Options")
    temperature = st.slider("LLM temperature", 0.0, 1.0, 0.0, 0.05)
    max_slides = st.number_input("Max slides", min_value=1, max_value=50, value=8)

st.markdown("Enter an outline or upload a DOCX, and the app will call Groq LLM to generate slides and a downloadable PPTX.")

outline = st.text_area("Input brief / outline", height=300)
uploaded = st.file_uploader("Upload DOCX (optional)", type=["docx"])  

if uploaded and Document:
    try:
        doc = Document(uploaded)
        text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        if not outline:
            outline = text
            st.success("Parsed uploaded DOCX into outline textarea")
    except Exception as e:
        st.warning(f"Failed to parse upload: {e}")

if st.button("Generate PPT"):
    if not outline.strip():
        st.error("Please provide an outline or upload a docx.")
    else:
        with st.spinner("Calling Groq LLM and building PPTX..."):
            slide_json = call_llm_for_slides(outline, temperature=temperature)
            if len(slide_json.get("slides", [])) > max_slides:
                slide_json["slides"] = slide_json["slides"][:max_slides]
            st.subheader("Generated slide JSON")
            st.json(slide_json)

            pptx_path = build_pptx_from_json(slide_json)
            st.success("PPTX generated")
            with open(pptx_path, "rb") as f:
                st.download_button("Download PPTX", f.read(), file_name=f"deck_{uuid.uuid4().hex[:6]}.pptx")

st.markdown("---")
st.caption("Prototype using Groq Cloud. For production: add auth, workers, storage, and integrations.")
